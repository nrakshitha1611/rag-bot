import os
import warnings
from pptx import Presentation
from langchain_community.vectorstores import Chroma
from langchain.embeddings.huggingface import HuggingFaceEmbeddings
from langchain.chains import RetrievalQA
from langchain.schema.runnable import RunnableLambda
from groq import Groq, APIError

# Suppress unnecessary warnings
warnings.filterwarnings("ignore", category=DeprecationWarning)
os.environ["TOKENIZERS_PARALLELISM"] = "false"

# Configuration constants
GROQ_API_KEY = "your_groq_api_key"
MAX_INPUT_TOKENS = 3000
MAX_OUTPUT_TOKENS = 500

def shorten_text(text, token_limit):
    """Trim text to meet token constraints."""
    if len(text.split()) > token_limit:
        return " ".join(text.split()[:token_limit]) + " [Text truncated]"
    return text

def fetch_ppt_content(file_path):
    """Extracts textual content from a PowerPoint file."""
    ppt = Presentation(file_path)
    all_text = []
    for slide in ppt.slides:
        for element in slide.shapes:
            if element.has_text_frame:
                all_text.append(element.text)
    return " ".join(all_text)

def initialize_vector_database(content):
    """Initialize ChromaDB with HuggingFace embeddings."""
    try:
        print("Setting up the vector database with embeddings...")
        embedder = HuggingFaceEmbeddings(model_name="sentence-transformers/all-mpnet-base-v2")
        vector_database = Chroma.from_texts(texts=[content], embedding=embedder)
        return vector_database
    except Exception as ex:
        print(f"Error while setting up the vector database: {ex}")
        raise

def configure_qa_system(vector_database):
    """Create a QA system using Groq API and the vector database."""
    retriever = vector_database.as_retriever(search_kwargs={"k": 1})

    def groq_response_engine(prompt, **kwargs):
        """Interact with Groq's LLM API."""
        try:
            if not isinstance(prompt, str):
                prompt = str(prompt)
            prompt = shorten_text(prompt, MAX_INPUT_TOKENS)

            print("Querying Groq's LLM for a response...")
            client = Groq(api_key=GROQ_API_KEY)
            response = client.chat.completions.create(
                messages=[{"role": "user", "content": prompt}],
                model="llama3-8b-8192",
                max_tokens=MAX_OUTPUT_TOKENS,
                **kwargs
            )

            # Extract the generated response correctly
            return response.choices[0].message.content
        except APIError as api_err:
            print(f"Groq API encountered an error: {api_err}")
            return "Unable to process your request. Please try again later."
        except Exception as gen_err:
            print(f"Unexpected error in Groq's LLM interaction: {gen_err}")
            return "An error occurred while generating a response."

    # Wrap the Groq engine in a Runnable
    qa_pipeline = RunnableLambda(groq_response_engine)
    return RetrievalQA.from_chain_type(llm=qa_pipeline, retriever=retriever)

def ask_questions():
    """Interact with the user for questions and answers."""
    print("Welcome to your AI-powered assistant!")
    file_path = input("Provide the PowerPoint file path to analyze: ").strip()
    if not os.path.exists(file_path):
        print("Oops! The file was not found. Please double-check the path.")
        return

    print("Extracting the content from your slides...")
    content = fetch_ppt_content(file_path)

    print("Building the database for content search...")
    try:
        vector_db = initialize_vector_database(content)
    except Exception as setup_err:
        print(f"Failed to initialize the system: {setup_err}")
        return

    print("Configuring the intelligent question-answer system...")
    try:
        qa_system = configure_qa_system(vector_db)
    except Exception as config_err:
        print(f"Failed to configure the QA system: {config_err}")
        return

    print("System is ready! Ask your questions below (type 'exit' to quit):")
    while True:
        user_input = input("\nYour query: ").strip()
        if user_input.lower() == "exit":
            print("Goodbye! Thanks for using this assistant.")
            break
        if not user_input:
            print("Please enter a valid question.")
            continue

        try:
            print("Thinking...")
            answer = qa_system.run(user_input)
            print(f"Answer: {answer}")
        except APIError as api_err:
            print(f"Encountered a problem while retrieving the response: {api_err}")
        except Exception as gen_err:
            print(f"Unexpected error: {gen_err}")

if __name__ == "__main__":
    ask_questions()
